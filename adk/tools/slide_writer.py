"""Slide writer — generates PowerPoint slides from analysis results.

Creates a structured presentation with title slide, insight slides
featuring charts, and a summary slide.
"""

import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTPUT_DIR = Path("output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)


def _add_background(slide, color: RGBColor):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_BLUE, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def write_slides(
    title: str,
    findings: list[str] | None = None,
    charts: list[str] | None = None,
    summary: str = "",
    model_used: str = "",
) -> Path:
    """Generate PowerPoint slides from analysis results.

    Creates:
    - Title slide
    - Summary slide
    - One slide per finding (with optional chart)
    - Summary/conclusion slide

    Args:
        title: Presentation title.
        findings: List of key findings (one per slide).
        charts: List of chart HTML file paths (one per finding).
        summary: Executive summary text.
        model_used: Model used for analysis.

    Returns:
        Path to the generated .pptx file.
    """
    findings = findings or []
    charts = charts or []

    # Sanitize title for filename
    safe_title = re.sub(r"[^\w\s-]", "", title).strip().lower()
    safe_title = re.sub(r"[\s]+", "-", safe_title)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    filename = f"{safe_title}-{timestamp}.pptx"
    output_path = OUTPUT_DIR / filename

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _add_background(slide, DARK_BLUE)
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(2),
                 title, font_size=36, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
                 f"Data Analyst Agent  |  {datetime.now().strftime('%Y-%m-%d %H:%M')}",
                 font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                 alignment=PP_ALIGN.CENTER)
    if model_used:
        _add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
                     f"Model: {model_used}", font_size=12,
                     color=RGBColor(0x88, 0x99, 0xAA),
                     alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Executive Summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 "EXECUTIVE SUMMARY", font_size=14, bold=True, color=ACCENT_BLUE)
    if summary:
        _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.5),
                     summary, font_size=16, color=DARK_BLUE)
    else:
        _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.5),
                     "No summary provided.", font_size=16, color=RGBColor(0x99, 0x99, 0x99))

    # --- Slides 3+: Findings ---
    for i, finding in enumerate(findings):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_background(slide, WHITE)

        # Finding header
        _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                     f"FINDING {i + 1}", font_size=14, bold=True, color=ACCENT_BLUE)

        # Finding text
        _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(7.5), Inches(5.5),
                     finding, font_size=16, color=DARK_BLUE)

        # Chart (if available for this finding)
        if i < len(charts):
            chart_path = charts[i]
            _add_textbox(slide, Inches(8.3), Inches(0.3), Inches(4.5), Inches(0.5),
                         "CHART", font_size=14, bold=True, color=ACCENT_BLUE,
                         alignment=PP_ALIGN.CENTER)
            _add_textbox(slide, Inches(8.3), Inches(1.0), Inches(4.5), Inches(0.3),
                         f"[{Path(chart_path).name}]", font_size=10,
                         color=RGBColor(0x88, 0x88, 0x88),
                         alignment=PP_ALIGN.CENTER)

    # --- Last Slide: Summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, DARK_BLUE)
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                 "SUMMARY", font_size=32, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(2),
                 f"Analysis completed using {model_used or 'default model'}.",
                 font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC),
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
                 font_size=12, color=RGBColor(0x88, 0x99, 0xAA),
                 alignment=PP_ALIGN.CENTER)

    prs.save(str(output_path))
    return output_path
